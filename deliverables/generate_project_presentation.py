#!/usr/bin/env python3
"""Generate the Policy Evidence Interface project update deck."""

from __future__ import annotations

import os
import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont

ROOT = Path(__file__).resolve().parent
ASSET_DIR = ROOT / "presentation_assets"
OUTPUT = ROOT / "Policy_Evidence_Interface_Generated_Draft.pptx"

W, H = 1600, 900

# Presentation-safe palette.
INK = "#111820"
INK_2 = "#25313D"
SLATE = "#5C6E81"
GOLD = "#E5B123"
GOLD_DARK = "#9A6A00"
CREAM = "#F6F3EA"
PAPER = "#FCFBF7"
MIST = "#E8EDF0"
WHITE = "#FFFFFF"
GREEN = "#2F6B57"
RED = "#9B3D36"

FONT_REG = "/usr/share/fonts/opentype/inter/Inter-Regular.otf"
FONT_MED = "/usr/share/fonts/opentype/inter/Inter-Medium.otf"
FONT_BOLD = "/usr/share/fonts/opentype/inter/Inter-Bold.otf"
FONT_DISPLAY = "/usr/share/fonts/opentype/inter/InterDisplay-SemiBold.otf"


def font(size: int, bold: bool = False, medium: bool = False, display: bool = False):
    path = FONT_DISPLAY if display else FONT_BOLD if bold else FONT_MED if medium else FONT_REG
    return ImageFont.truetype(path, size=size)


def rounded(draw: ImageDraw.ImageDraw, box, radius=22, fill=WHITE, outline=None, width=2):
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)


def line(draw: ImageDraw.ImageDraw, xy, fill=SLATE, width=4):
    draw.line(xy, fill=fill, width=width)


def wrap(draw: ImageDraw.ImageDraw, text: str, fnt, max_width: int) -> list[str]:
    if not text:
        return [""]
    out: list[str] = []
    for paragraph in text.split("\n"):
        if paragraph == "":
            out.append("")
            continue
        words = paragraph.split()
        current = words[0]
        for word in words[1:]:
            trial = f"{current} {word}"
            if draw.textlength(trial, font=fnt) <= max_width:
                current = trial
            else:
                out.append(current)
                current = word
        out.append(current)
    return out


def text_block(
    draw: ImageDraw.ImageDraw,
    xy,
    text: str,
    fnt,
    fill=INK,
    max_width=1000,
    spacing=8,
    anchor=None,
):
    x, y = xy
    lines = wrap(draw, text, fnt, max_width)
    ascent, descent = fnt.getmetrics()
    line_height = ascent + descent + spacing
    for i, txt in enumerate(lines):
        draw.text((x, y + i * line_height), txt, font=fnt, fill=fill, anchor=anchor)
    return y + len(lines) * line_height


def title(draw: ImageDraw.ImageDraw, kicker: str, heading: str, sub: str | None = None):
    draw.rectangle((74, 58, 84, 115), fill=GOLD)
    draw.text((108, 58), kicker.upper(), font=font(19, bold=True), fill=GOLD_DARK)
    draw.text((108, 87), heading, font=font(42, display=True), fill=INK)
    if sub:
        draw.text((108, 145), sub, font=font(21), fill=SLATE)


def footer(draw: ImageDraw.ImageDraw, page: int, source: str = ""):
    draw.text((1526, 853), f"{page:02d}", font=font(15, bold=True), fill=INK, anchor="ra")


def pill(draw, box, label, fill=INK, text_fill=WHITE):
    rounded(draw, box, radius=20, fill=fill, outline=None)
    x1, y1, x2, y2 = box
    draw.text(((x1 + x2) / 2, (y1 + y2) / 2), label, font=font(16, bold=True), fill=text_fill, anchor="mm")


def icon_nodes(draw: ImageDraw.ImageDraw, cx: int, cy: int):
    pts = [(cx, cy), (cx + 185, cy - 100), (cx + 215, cy + 100), (cx + 395, cy - 25)]
    for a, b in [(0, 1), (0, 2), (1, 3), (2, 3), (1, 2)]:
        line(draw, (*pts[a], *pts[b]), fill=SLATE, width=5)
    sizes = [26, 35, 30, 45]
    colors = [GOLD, CREAM, WHITE, GOLD]
    for (x, y), r, c in zip(pts, sizes, colors):
        draw.ellipse((x - r, y - r, x + r, y + r), fill=c, outline=INK, width=5)
    draw.ellipse((cx + 365, cy - 55, cx + 425, cy + 5), fill=INK)
    draw.text((cx + 395, cy - 25), "AI", font=font(19, bold=True), fill=WHITE, anchor="mm")


def base(bg=PAPER):
    return Image.new("RGB", (W, H), bg)


def slide_1() -> Image.Image:
    im = base(INK)
    d = ImageDraw.Draw(im)
    d.rectangle((0, 0, 26, H), fill=GOLD)
    d.text((90, 74), "POLICY EVIDENCE INTERFACE", font=font(20, bold=True), fill=GOLD)
    d.text((90, 154), "Technical status of the", font=font(54, display=True), fill=WHITE)
    d.text((90, 220), "Drupal MCP module", font=font(54, display=True), fill=WHITE)
    text_block(
        d,
        (92, 318),
        "Codebase review • implemented interfaces • engineering backlog",
        font(24),
        fill="#D9DFE3",
        max_width=700,
        spacing=13,
    )
    rounded(d, (90, 420, 790, 565), radius=24, fill=GOLD)
    d.text((122, 445), "CURRENT STATUS", font=font(17, bold=True), fill=INK)
    text_block(
        d,
        (122, 480),
        "The module can serve published Drupal content and text-extractable PDF pages to MCP clients.",
        font(25, medium=True),
        fill=INK,
        max_width=630,
        spacing=5,
    )
    d.text((90, 603), "WORKING PROOF OF CONCEPT", font=font(17, bold=True), fill=GREEN)
    d.text((352, 603), "•", font=font(17, bold=True), fill=SLATE)
    d.text((380, 603), "NOT YET PRODUCTION READY", font=font(17, bold=True), fill="#D8877F")

    stat_cards = [
        (890, 168, "5", "MCP TOOLS"),
        (1190, 168, "3", "RESOURCES"),
        (890, 365, "2", "TRANSPORTS"),
        (1190, 365, "2024-11-05", "MCP PROTOCOL"),
    ]
    for x, y, value, label in stat_cards:
        rounded(d, (x, y, x + 250, y + 155), radius=20, fill=INK_2, outline="#43515D", width=2)
        size = 36 if len(value) > 3 else 52
        d.text((x + 28, y + 28), value, font=font(size, display=True), fill=GOLD)
        d.text((x + 28, y + 105), label, font=font(16, bold=True), fill=WHITE)

    rounded(d, (890, 578, 1440, 662), radius=18, fill="#202B35", outline="#43515D", width=2)
    d.text((918, 600), "NEXT CODE MILESTONE", font=font(15, bold=True), fill=GOLD)
    d.text((918, 628), "Secure, test and deploy on the APO test site", font=font(18, medium=True), fill=WHITE)

    d.text((90, 760), "Current branch: main • 27 commits • PHP/Drupal module", font=font(18, medium=True), fill=WHITE)
    d.text((90, 800), "Project update • 31 July 2026", font=font(17), fill="#AAB4BC")
    d.text((1510, 800), "01", font=font(17, bold=True), fill=GOLD, anchor="ra")
    return im


def slide_2() -> Image.Image:
    im = base()
    d = ImageDraw.Draw(im)
    title(d, "01 • architecture", "How the module processes an MCP request", "Two transports share the same plugin-based tool and resource surface.")

    boxes = [
        (74, 222, 302, 364, "CLIENT", "MCP-capable AI app", False),
        (354, 222, 620, 364, "TRANSPORT", "POST /_mcp\nor drush mcp:server", True),
        (672, 222, 974, 364, "DISPATCH", "Controller\nor Commands class", False),
        (1026, 222, 1252, 364, "DISCOVERY", "plugin managers", False),
        (1304, 222, 1526, 364, "EXECUTE", "tool / resource", False),
    ]
    for i, (x1, y1, x2, y2, head, body, accent) in enumerate(boxes):
        rounded(d, (x1, y1, x2, y2), radius=20, fill=INK if accent else WHITE, outline=None if accent else "#C9D2D8", width=2)
        d.text(((x1 + x2) / 2, y1 + 36), head, font=font(18, bold=True), fill=GOLD if accent else INK, anchor="mm")
        lines = body.split("\n")
        for li, txt in enumerate(lines):
            d.text(((x1 + x2) / 2, y1 + 82 + li * 26), txt, font=font(16), fill=WHITE if accent else SLATE, anchor="mm")
        if i < len(boxes) - 1:
            arrow(d, x2 + 8, boxes[i + 1][0] - 10, 293)

    rounded(d, (74, 414, 742, 748), radius=22, fill=WHITE, outline="#D3DADF", width=2)
    d.text((108, 448), "REQUEST LIFECYCLE", font=font(18, bold=True), fill=GREEN)
    lifecycle = [
        ("1", "JSON-RPC method", "initialize, tools/list, tools/call, resources/list or resources/read"),
        ("2", "Dispatcher", "validates the method and routes to a handler"),
        ("3", "Plugin manager", "discovers annotation-defined classes and instantiates the match"),
        ("4", "Response", "returns JSON-RPC result or a structured protocol error"),
    ]
    y = 490
    for num, head, body in lifecycle:
        d.ellipse((108, y, 146, y + 38), fill=MIST)
        d.text((127, y + 19), num, font=font(16, bold=True), fill=INK, anchor="mm")
        d.text((164, y), head, font=font(17, bold=True), fill=INK)
        text_block(d, (164, y + 26), body, font(15), fill=SLATE, max_width=530, spacing=3)
        y += 65

    rounded(d, (780, 414, 1526, 748), radius=22, fill=INK)
    d.text((814, 448), "CODE MAP", font=font(18, bold=True), fill=GOLD)
    code_map = [
        ("HTTP", "src/Controller/McpServerController.php"),
        ("STDIO", "src/Commands/McpServerCommands.php"),
        ("PLUGINS", "src/Plugin/McpTool/* + McpResource/*"),
        ("DATA", "Drupal Entity API + smalot/pdfparser"),
    ]
    y = 496
    for label, path in code_map:
        pill(d, (814, y, 930, y + 36), label, fill="#354451", text_fill=GOLD)
        d.text((952, y + 18), path, font=font(15, medium=True), fill=WHITE, anchor="lm")
        y += 48
    rounded(d, (814, 690, 1490, 728), radius=14, fill=GOLD)
    d.text((832, 703), "Extension model: add one annotated plugin class; discovery is automatic.", font=font(16, bold=True), fill=INK, anchor="lm")
    footer(d, 2, "Architecture verified from controllers, Drush command, plugin managers and annotated plugin classes")
    return im


def arch_box(draw, box, top, bottom, accent=False):
    fill = INK if accent else WHITE
    outline = None if accent else "#CBD3D8"
    rounded(draw, box, radius=18, fill=fill, outline=outline, width=2)
    x1, y1, x2, y2 = box
    draw.text(((x1 + x2) / 2, y1 + 36), top, font=font(20, bold=True), fill=GOLD if accent else INK, anchor="mm")
    draw.text(((x1 + x2) / 2, y1 + 77), bottom, font=font(16), fill=WHITE if accent else SLATE, anchor="mm")


def arrow(draw, x1, x2, y):
    line(draw, (x1, y, x2 - 10, y), fill=GOLD_DARK, width=5)
    draw.polygon([(x2, y), (x2 - 18, y - 10), (x2 - 18, y + 10)], fill=GOLD_DARK)


def bullet_list(draw, x, y, bullets, color=INK, mark=GOLD, max_width=600, size=18, gap=48):
    for item in bullets:
        draw.ellipse((x, y + 7, x + 12, y + 19), fill=mark)
        text_block(draw, (x + 26, y), item, font(size), fill=color, max_width=max_width, spacing=4)
        y += gap
    return y


def slide_3() -> Image.Image:
    im = base()
    d = ImageDraw.Draw(im)
    title(d, "02 • features", "Implemented MCP surface", "Five tools expose site structure, published node data and page-level PDF text.")

    rounded(d, (74, 214, 960, 748), radius=22, fill=WHITE, outline="#D3DADF", width=2)
    d.text((108, 246), "TOOLS", font=font(18, bold=True), fill=GREEN)
    d.text((360, 246), "INPUT", font=font(15, bold=True), fill=SLATE)
    d.text((610, 246), "CURRENT BEHAVIOUR", font=font(15, bold=True), fill=SLATE)
    line(d, (108, 278, 926, 278), fill="#D9DEE2", width=2)
    tools = [
        ("get_site_info", "none", "Site name, URL, Drupal version and language"),
        ("list_content_types", "none", "Configured node bundles and descriptions"),
        ("search_nodes", "keyword, type?, limit?", "Title LIKE search; published results; limit ≤ 50"),
        ("get_node", "nid", "Published node body, author, URL and field values"),
        ("read_node_pdf", "nid, page_start?", "One validated page from an attached PDF"),
    ]
    y = 295
    for i, (name, inp, behaviour) in enumerate(tools):
        if i % 2 == 0:
            rounded(d, (98, y - 8, 936, y + 66), radius=12, fill="#F5F7F8")
        d.text((112, y + 16), name, font=font(17, bold=True), fill=INK, anchor="lm")
        d.text((360, y + 16), inp, font=font(15, medium=True), fill=GOLD_DARK, anchor="lm")
        text_block(d, (610, y), behaviour, font(15), fill=INK_2, max_width=300, spacing=3)
        y += 83

    rounded(d, (998, 214, 1526, 494), radius=22, fill=INK)
    d.text((1032, 246), "RESOURCES", font=font(18, bold=True), fill=GOLD)
    resources = [
        ("drupal://site-info", "site configuration"),
        ("drupal://content-types", "configured bundles"),
        ("drupal://recent-nodes", "20 latest published nodes"),
    ]
    y = 292
    for uri, body in resources:
        rounded(d, (1032, y, 1492, y + 54), radius=12, fill="#263440")
        d.text((1050, y + 16), uri, font=font(15, bold=True), fill=WHITE)
        d.text((1474, y + 17), body, font=font(13), fill="#BFC8CE", anchor="ra")
        y += 67

    rounded(d, (998, 520, 1526, 748), radius=22, fill="#F1E7CA")
    d.text((1032, 552), "SUPPORTED JSON-RPC METHODS", font=font(17, bold=True), fill=GOLD_DARK)
    methods = [
        "initialize  •  ping",
        "tools/list  •  tools/call",
        "resources/list  •  resources/read",
        "resources/templates/list  •  prompts/list",
    ]
    y = 596
    for method in methods:
        d.ellipse((1034, y + 6, 1046, y + 18), fill=GOLD_DARK)
        d.text((1060, y), method, font=font(15, medium=True), fill=INK)
        y += 36

    rounded(d, (74, 772, 1526, 820), radius=14, fill=INK)
    d.text((100, 796), "CONTENT RULE", font=font(15, bold=True), fill=GOLD, anchor="lm")
    d.text((262, 796), "Node retrieval checks publication status; search and recent-node queries also use Drupal access checks.", font=font(17, medium=True), fill=WHITE, anchor="lm")
    footer(d, 3, "Feature inventory verified from src/Plugin/McpTool and src/Plugin/McpResource")
    return im


def timeline_stage(draw, x, w, num, sprint, date, heading, items, fill):
    rounded(draw, (x, 232, x + w, 742), radius=22, fill=WHITE, outline="#D3DADF", width=2)
    draw.rectangle((x, 232, x + w, 245), fill=fill)
    draw.ellipse((x + 28, 270, x + 82, 324), fill=fill)
    draw.text((x + 55, 297), str(num), font=font(21, bold=True), fill=INK if fill == GOLD else WHITE, anchor="mm")
    draw.text((x + 102, 272), sprint, font=font(17, bold=True), fill=INK)
    draw.text((x + 102, 300), date, font=font(15, medium=True), fill=SLATE)
    text_block(draw, (x + 28, 354), heading, font(25, display=True), fill=INK, max_width=w - 56, spacing=5)
    bullet_list(draw, x + 30, 430, items, color=INK_2, mark=fill, max_width=w - 82, size=17, gap=66)


def slide_4() -> Image.Image:
    im = base()
    d = ImageDraw.Draw(im)
    title(d, "03 • engineering review", "Code strengths and production gaps", "The proof of concept is intentionally small; the missing controls are now specific and actionable.")

    rounded(d, (74, 218, 780, 742), radius=22, fill=WHITE, outline="#D1DADF", width=2)
    d.text((108, 252), "ALREADY PRESENT IN THE CODE", font=font(18, bold=True), fill=GREEN)
    bullet_list(
        d,
        110,
        302,
        [
            "Restricted Drupal permission protects the HTTP endpoint",
            "Node tools reject missing and unpublished content",
            "Search has Drupal access checks and a hard result cap",
            "PDF tool validates files, readability and page range",
            "JSON-RPC parse, method and internal errors are returned",
            "Annotation-based plugins keep extensions modular",
            "CI validates Composer metadata and PHP syntax",
        ],
        color=INK_2,
        mark=GREEN,
        max_width=610,
        size=17,
        gap=57,
    )

    rounded(d, (818, 218, 1526, 742), radius=22, fill="#F8ECE9", outline="#E3C3BD", width=2)
    d.text((852, 252), "GAPS TO CLOSE BEFORE DEPLOYMENT", font=font(18, bold=True), fill=RED)
    bullet_list(
        d,
        854,
        302,
        [
            "Wildcard CORS; no origin allow-list or session policy",
            "One Drupal permission; no client or licence-tier model",
            "Direct node/PDF reads lack per-node access checks",
            "No rate limiting, structured audit log or monitoring",
            "Title-only LIKE search; no semantic index or ranking",
            "No OCR or cached extraction for difficult PDFs",
            "No automated unit, integration or end-to-end test suite",
            "Responses expose admin email and absolute file path",
        ],
        color=INK_2,
        mark=RED,
        max_width=610,
        size=17,
        gap=50,
    )

    rounded(d, (74, 770, 1526, 820), radius=14, fill=INK)
    d.text((100, 796), "FIRST HARDENING PASS", font=font(15, bold=True), fill=GOLD, anchor="lm")
    d.text((332, 796), "Least privilege • response minimisation • origin/rate controls • auditability • automated tests", font=font(17, medium=True), fill=WHITE, anchor="lm")
    footer(d, 4, "Engineering review based on routing, controllers, tools, resources, CI and current response payloads")
    return im


def decision_card(draw, box, num, heading, body):
    x1, y1, x2, y2 = box
    rounded(draw, box, radius=20, fill=WHITE, outline="#D1D9DD", width=2)
    draw.ellipse((x1 + 24, y1 + 24, x1 + 74, y1 + 74), fill=GOLD)
    draw.text((x1 + 49, y1 + 49), str(num), font=font(19, bold=True), fill=INK, anchor="mm")
    draw.text((x1 + 92, y1 + 28), heading, font=font(20, bold=True), fill=INK)
    text_block(draw, (x1 + 92, y1 + 62), body, font(17), fill=SLATE, max_width=x2 - x1 - 120, spacing=4)


def slide_5() -> Image.Image:
    im = base()
    d = ImageDraw.Draw(im)
    title(d, "04 • implementation plan", "Codebase backlog for Sprints 4–6", "The sequence moves from secure deployment, to retrieval quality, to release assurance.")

    stage_w = 452
    timeline_stage(
        d, 74, stage_w, 1, "SPRINT 4", "review • 16 Aug",
        "Harden the MCP core",
        [
            "Deploy current module to the APO test site",
            "Implement client access model and minimise responses",
            "Restrict CORS; add rate limits and structured logging",
            "Test initialize, list, call, read and error paths",
        ],
        GOLD,
    )
    timeline_stage(
        d, 574, stage_w, 2, "SPRINT 5", "review • 6 Sep",
        "Build document retrieval",
        [
            "Extract PDF handling into a reusable service",
            "Add OCR fallback and cached text extraction",
            "Preserve attachment metadata and citation provenance",
            "Prototype semantic search; decide MCP/RAG boundary",
        ],
        SLATE,
    )
    timeline_stage(
        d, 1074, stage_w, 3, "SPRINT 6", "review • 11 Oct",
        "Validate and package",
        [
            "Measure relevance, latency and failure behaviour",
            "Complete security review and fail-safe testing",
            "Finish configuration, deployment and operator docs",
            "Package release candidate and handover evidence",
        ],
        GREEN,
    )

    rounded(d, (74, 770, 1526, 822), radius=16, fill=INK)
    d.text((100, 796), "DEFINITION OF DONE", font=font(15, bold=True), fill=GOLD, anchor="lm")
    d.text((330, 796), "Representative APO content works end-to-end with controlled access, observable failures and reproducible tests.", font=font(17, medium=True), fill=WHITE, anchor="lm")
    footer(d, 5, "Implementation sequence derived from code gaps, sprint dates and 31 July technical priorities")
    return im


def simple_slide_1() -> Image.Image:
    """Refined cover slide."""
    im = base(INK)
    d = ImageDraw.Draw(im)
    d.rectangle((0, 0, 24, H), fill=GOLD)
    d.text((100, 92), "TECHNICAL PROJECT UPDATE", font=font(19, bold=True), fill=GOLD)
    d.text((100, 278), "Policy Evidence", font=font(64, display=True), fill=WHITE)
    d.text((100, 354), "Interface", font=font(64, display=True), fill=WHITE)
    d.rectangle((102, 470, 350, 478), fill=GOLD)
    d.text((100, 520), "A controlled AI interface for Drupal policy evidence", font=font(26, medium=True), fill="#D5DDE2")
    d.text((100, 790), "ANU TechLauncher  •  31 July 2026", font=font(16), fill="#89949D")

    rounded(d, (980, 112, 1492, 734), radius=24, fill="#1B2731", outline="#354451", width=2)
    d.text((1020, 152), "PROJECT TEAM", font=font(16, bold=True), fill=GOLD)
    team = ["Harry Baard", "Amogh Agarwal", "Peter Wei", "Qiyue Zhang", "Shiyun Yao"]
    y = 202
    for member in team:
        d.text((1020, y), member, font=font(19, medium=True), fill=WHITE)
        y += 42
    line(d, (1020, 432, 1452, 432), fill="#42515D", width=2)
    d.text((1020, 468), "CURRENT TUTOR", font=font(14, bold=True), fill=GOLD)
    d.text((1020, 500), "Zara Hassan", font=font(18), fill="#D4DCE1")
    d.text((1020, 558), "CLIENT", font=font(14, bold=True), fill=GOLD)
    d.text((1020, 590), "Dr. Harry Rolf", font=font(18), fill="#D4DCE1")
    d.text((1020, 622), "Shape Policy", font=font(16), fill="#9FABB3")
    d.text((1510, 800), "01", font=font(17, bold=True), fill=GOLD, anchor="ra")
    return im


def simple_slide_reason() -> Image.Image:
    """Project vision and design principles."""
    im = base()
    d = ImageDraw.Draw(im)
    title(d, "01 • vision", "Design intent", "A controlled machine interface for curated policy evidence.")

    rounded(d, (74, 220, 570, 706), radius=24, fill=INK)
    d.text((112, 260), "THE TENSION", font=font(17, bold=True), fill=GOLD)
    d.text((112, 336), "AI NEEDS", font=font(15, bold=True), fill="#AAB4BC")
    text_block(d, (112, 374), "Fast, direct access.", font(30, display=True), fill=WHITE, max_width=400, spacing=7)
    line(d, (112, 456, 524, 456), fill="#42515D", width=2)
    d.text((112, 496), "EVIDENCE NEEDS", font=font(15, bold=True), fill="#AAB4BC")
    text_block(d, (112, 534), "Provenance, permissions and context.", font(30, display=True), fill=WHITE, max_width=400, spacing=7)

    principles = [
        (620, 220, GOLD, "1", "GROUNDED", "Answers begin with published Drupal content and retain source links."),
        (620, 378, SLATE, "2", "GOVERNED", "Access and future licensing rules remain under repository control."),
        (620, 536, GREEN, "3", "EXTENSIBLE", "New retrieval functions can be added through Drupal plugins."),
    ]
    for x, y, color, num, heading, body in principles:
        rounded(d, (x, y, 1526, y + 132), radius=20, fill=WHITE, outline="#D1DADF", width=2)
        d.ellipse((x + 28, y + 33, x + 94, y + 99), fill=color)
        d.text((x + 61, y + 66), num, font=font(24, bold=True), fill=INK if color == GOLD else WHITE, anchor="mm")
        d.text((x + 122, y + 26), heading, font=font(19, bold=True), fill=INK)
        text_block(d, (x + 122, y + 61), body, font(18), fill=SLATE, max_width=735, spacing=5)

    rounded(d, (74, 744, 1526, 814), radius=16, fill="#F1E7CA")
    d.text((100, 779), "DESIGN PRINCIPLE", font=font(14, bold=True), fill=GOLD_DARK, anchor="lm")
    d.text((300, 779), "Drupal remains the source of truth; MCP is the access layer, not a replacement repository.", font=font(18, medium=True), fill=INK, anchor="lm")
    footer(d, 2)
    return im


def simple_slide_2() -> Image.Image:
    """End-to-end proof provided by the current prototype."""
    im = base()
    d = ImageDraw.Draw(im)
    title(d, "02 • prototype", "Prototype proof", "What the current code demonstrates end to end.")

    stages = [
        (74, 220, 230, "AI CLIENT", "MCP request"),
        (348, 220, 250, "TRANSPORT", "HTTP or stdio"),
        (642, 220, 260, "DISPATCH", "JSON-RPC methods"),
        (946, 220, 260, "PLUGIN LAYER", "tools + resources"),
        (1250, 220, 276, "DRUPAL DATA", "nodes + PDFs"),
    ]
    for i, (x, y, w, heading, body) in enumerate(stages):
        fill = INK if i == 2 else WHITE
        rounded(d, (x, y, x + w, y + 122), radius=18, fill=fill, outline=None if fill == INK else "#C9D2D8", width=2)
        d.text((x + w / 2, y + 39), heading, font=font(17, bold=True), fill=GOLD if fill == INK else INK, anchor="mm")
        d.text((x + w / 2, y + 82), body, font=font(15), fill=WHITE if fill == INK else SLATE, anchor="mm")
        if i < len(stages) - 1:
            arrow(d, x + w + 10, stages[i + 1][0] - 10, y + 61)

    capabilities = [
        (74, GOLD, "CONNECT", "2 transports", "Browser/server clients and local command-line clients."),
        (448, SLATE, "DISCOVER", "5 tools • 3 resources", "Clients can list capabilities before calling them."),
        (822, GREEN, "RETRIEVE", "nodes • fields • PDF page", "Published content is returned as structured JSON text."),
        (1196, INK_2, "EXTEND", "annotation plugins", "New functions are discovered through Drupal’s plugin system."),
    ]
    for x, color, heading, metric, body in capabilities:
        rounded(d, (x, 400, x + 330, 704), radius=22, fill=WHITE, outline="#D1DADF", width=2)
        d.rectangle((x, 400, x + 330, 412), fill=color)
        d.text((x + 28, 444), heading, font=font(17, bold=True), fill=INK)
        d.text((x + 28, 486), metric, font=font(24, display=True), fill=color)
        text_block(d, (x + 28, 548), body, font(17), fill=SLATE, max_width=270, spacing=5)

    rounded(d, (74, 744, 1526, 814), radius=16, fill=INK)
    d.text((100, 779), "PROOF POINT", font=font(14, bold=True), fill=GOLD, anchor="lm")
    d.text((262, 779), "A standard MCP client can discover and invoke Drupal functions without a repository-specific client.", font=font(18, medium=True), fill=WHITE, anchor="lm")
    footer(d, 3)
    return im


def simple_slide_3() -> Image.Image:
    """Code-informed readiness assessment."""
    im = base()
    d = ImageDraw.Draw(im)
    title(d, "03 • assessment", "Readiness gaps", "The prototype is extensible; access, operations and retrieval depth are not yet production-ready.")

    columns = [
        (74, RED, "ACCESS & EXPOSURE", [
            "Direct node/PDF reads do not enforce entity view access.",
            "One route permission covers every client and content type.",
            "Wildcard origins and responses expose email/file-path details.",
        ], "Risk: published but restricted content may be overexposed."),
        (574, GOLD, "RELIABILITY & OPERATIONS", [
            "HTTP and stdio duplicate the protocol dispatcher.",
            "A faulty plugin can disrupt tool/resource discovery.",
            "Logging/templates are partial; no behavioural test suite.",
        ], "Risk: maintenance drift and fragile client connections."),
        (1074, SLATE, "RETRIEVAL QUALITY", [
            "Search matches titles only; there is no ranking or semantic index.",
            "PDF extraction returns one text-based page at a time.",
            "No OCR, cached extraction or attachment-level search.",
        ], "Risk: the demo works, but evidence coverage is narrow."),
    ]
    for x, color, heading, items, implication in columns:
        rounded(d, (x, 220, x + 452, 714), radius=22, fill=WHITE, outline="#D1DADF", width=2)
        d.rectangle((x, 220, x + 452, 234), fill=color)
        d.text((x + 28, 264), heading, font=font(17, bold=True), fill=INK)
        y = 316
        for item in items:
            d.ellipse((x + 30, y + 8, x + 43, y + 21), fill=color)
            text_block(d, (x + 58, y), item, font(16), fill=INK_2, max_width=350, spacing=4)
            y += 94
        rounded(d, (x + 24, 610, x + 428, 686), radius=14, fill="#F4F6F7")
        text_block(d, (x + 42, 630), implication, font(15, medium=True), fill=INK, max_width=365, spacing=4)

    rounded(d, (74, 744, 1526, 814), radius=16, fill=INK)
    d.text((100, 779), "READINESS", font=font(14, bold=True), fill=GOLD, anchor="lm")
    d.text((244, 779), "Appropriate for controlled internal testing only; access hardening is the release gate.", font=font(18, medium=True), fill=WHITE, anchor="lm")
    footer(d, 4)
    return im


def simple_slide_4() -> Image.Image:
    """Decision-oriented pilot blueprint."""
    im = base()
    d = ImageDraw.Draw(im)
    title(d, "04 • delivery", "Pilot blueprint", "Sequence the work around secure access, retrieval quality and operational assurance.")

    workstreams = [
        (74, GOLD, "SPRINT 4", "HARDEN ACCESS", [
            "Share one dispatcher and response model across both transports.",
            "Apply entity-level permissions and minimise response data.",
            "Add origin controls, usage limits and audit logging.",
        ], "SUCCESS", "Unauthorised content is never returned."),
        (574, SLATE, "SPRINT 5", "IMPROVE RETRIEVAL", [
            "Create a reusable extraction service with OCR and caching.",
            "Establish a search baseline, then prototype RAG.",
            "Return citation metadata and stable source links.",
        ], "SUCCESS", "Representative nodes and PDFs are retrievable with provenance."),
        (1074, GREEN, "SPRINT 6", "PROVE OPERATION", [
            "Test the same behaviours over HTTP and command line.",
            "Isolate plugin failures and make errors observable.",
            "Deploy to the test site and complete security review.",
        ], "SUCCESS", "The demonstration is repeatable and failures are diagnosable."),
    ]
    for x, color, sprint, heading, items, success_label, success in workstreams:
        rounded(d, (x, 220, x + 452, 714), radius=22, fill=WHITE, outline="#D1DADF", width=2)
        d.rectangle((x, 220, x + 452, 234), fill=color)
        pill(d, (x + 28, 260, x + 142, 296), sprint, fill="#EEF1F3", text_fill=INK)
        d.text((x + 28, 322), heading, font=font(19, bold=True), fill=INK)
        y = 372
        for item in items:
            d.ellipse((x + 30, y + 8, x + 43, y + 21), fill=color)
            text_block(d, (x + 58, y), item, font(15), fill=INK_2, max_width=350, spacing=4)
            y += 78
        rounded(d, (x + 24, 612, x + 428, 686), radius=14, fill="#F4F6F7")
        d.text((x + 42, 626), success_label, font=font(12, bold=True), fill=color)
        text_block(d, (x + 42, 650), success, font(14, medium=True), fill=INK, max_width=360, spacing=3)

    rounded(d, (74, 744, 1526, 814), radius=16, fill=INK)
    d.text((100, 779), "FIRST DECISION", font=font(14, bold=True), fill=GOLD, anchor="lm")
    d.text((284, 779), "Agree the pilot users and representative content set before expanding search features.", font=font(18, medium=True), fill=WHITE, anchor="lm")
    footer(d, 5)
    return im


def make_contact_sheet(paths: list[Path]):
    thumbs = []
    for p in paths:
        img = Image.open(p)
        img.thumbnail((640, 360))
        thumbs.append(img.copy())
    rows = (len(thumbs) + 1) // 2
    sheet = Image.new("RGB", (1320, rows * 370 + 20), "#DCE2E5")
    coords = [(20 + (i % 2) * 660, 20 + (i // 2) * 370) for i in range(len(thumbs))]
    sd = ImageDraw.Draw(sheet)
    for i, (thumb, (x, y)) in enumerate(zip(thumbs, coords), 1):
        sheet.paste(thumb, (x, y))
        sd.text((x + 8, y + 326), f"SLIDE {i}", font=font(14, bold=True), fill=WHITE)
    sheet.save(ASSET_DIR / "contact_sheet.png", quality=94)


def build_pptx(paths: list[Path]):
    sys.path.insert(0, "/tmp/policy_evidence_pptx_lib")
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    def rgb(value: str):
        return RGBColor.from_string(value.lstrip("#"))

    def pos(value: float):
        return Inches(value / 120)

    def add_box(slide, x, y, w, h, fill, outline=None, rounded_box=True):
        kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded_box else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(kind, pos(x), pos(y), pos(w), pos(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        if outline:
            shape.line.color.rgb = rgb(outline)
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        return shape

    def add_circle(slide, x, y, size, fill):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, pos(x), pos(y), pos(size), pos(size))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        shape.line.fill.background()
        return shape

    def add_text(
        slide,
        x,
        y,
        w,
        h,
        value,
        size=18,
        color=INK,
        bold=False,
        align="left",
        valign="top",
        margin=0,
        font_name="Arial",
    ):
        shape = slide.shapes.add_textbox(pos(x), pos(y), pos(w), pos(h))
        shape.name = f"Text: {value[:35]}"
        frame = shape.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = pos(margin)
        frame.margin_right = pos(margin)
        frame.margin_top = pos(margin)
        frame.margin_bottom = pos(margin)
        frame.vertical_anchor = {
            "top": MSO_ANCHOR.TOP,
            "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }[valign]
        paragraph = frame.paragraphs[0]
        paragraph.text = value
        paragraph.alignment = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }[align]
        paragraph.font.name = font_name
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = rgb(color)
        paragraph.space_after = Pt(0)
        return shape

    def add_arrow(slide, x, y, w):
        arrow_shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, pos(x), pos(y), pos(w), pos(24))
        arrow_shape.fill.solid()
        arrow_shape.fill.fore_color.rgb = rgb(GOLD_DARK)
        arrow_shape.line.fill.background()
        return arrow_shape

    def add_title(slide, section, heading, subtitle):
        add_box(slide, 74, 58, 10, 57, GOLD, rounded_box=False)
        add_text(slide, 108, 56, 650, 30, section.upper(), 14, GOLD_DARK, True)
        add_text(slide, 108, 88, 1350, 54, heading, 31, INK, True)
        add_text(slide, 108, 145, 1350, 34, subtitle, 16, SLATE)

    def add_footer(slide, page, source):
        add_text(slide, 1460, 850, 66, 24, f"{page:02d}", 11, INK, True, "right")

    def add_bullet(slide, x, y, text, color, width=590, size=17):
        add_circle(slide, x, y + 8, 12, color)
        add_text(slide, x + 26, y, width, 46, text, size, INK_2, False, "left", "top")

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1 — refined cover and people.
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(INK)
    add_box(slide, 0, 0, 24, 900, GOLD, rounded_box=False)
    add_text(slide, 100, 92, 600, 34, "TECHNICAL PROJECT UPDATE", 14, GOLD, True)
    add_text(slide, 100, 278, 820, 150, "Policy Evidence\nInterface", 43, WHITE, True)
    add_box(slide, 102, 470, 248, 8, GOLD, rounded_box=False)
    add_text(slide, 100, 520, 790, 72, "A controlled AI interface for Drupal policy evidence", 19, "#D5DDE2")
    add_text(slide, 100, 790, 500, 24, "ANU TechLauncher  •  31 July 2026", 11, "#89949D")
    add_box(slide, 980, 112, 512, 622, "#1B2731", "#354451")
    add_text(slide, 1020, 152, 300, 26, "PROJECT TEAM", 12, GOLD, True)
    for i, member in enumerate(["Harry Baard", "Amogh Agarwal", "Peter Wei", "Qiyue Zhang", "Shiyun Yao"]):
        add_text(slide, 1020, 200 + i * 42, 390, 28, member, 14, WHITE, True)
    add_box(slide, 1020, 432, 432, 2, "#42515D", rounded_box=False)
    add_text(slide, 1020, 468, 220, 24, "CURRENT TUTOR", 10, GOLD, True)
    add_text(slide, 1020, 500, 300, 28, "Zara Hassan", 13, "#D4DCE1")
    add_text(slide, 1020, 558, 180, 24, "CLIENT", 10, GOLD, True)
    add_text(slide, 1020, 590, 300, 28, "Dr. Harry Rolf", 13, "#D4DCE1")
    add_text(slide, 1020, 622, 300, 24, "Shape Policy", 11, "#9FABB3")
    add_text(slide, 1460, 790, 50, 24, "01", 11, GOLD, True, "right")

    # Slide 2 — design intent.
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(PAPER)
    add_title(slide, "01 • vision", "Design intent", "A controlled machine interface for curated policy evidence.")
    add_box(slide, 74, 220, 496, 486, INK)
    add_text(slide, 112, 260, 240, 26, "THE TENSION", 12, GOLD, True)
    add_text(slide, 112, 336, 180, 24, "AI NEEDS", 10, "#AAB4BC", True)
    add_text(slide, 112, 374, 400, 50, "Fast, direct access.", 22, WHITE, True)
    add_box(slide, 112, 456, 412, 2, "#42515D", rounded_box=False)
    add_text(slide, 112, 496, 220, 24, "EVIDENCE NEEDS", 10, "#AAB4BC", True)
    add_text(slide, 112, 534, 400, 90, "Provenance, permissions\nand context.", 22, WHITE, True)
    principles = [
        (620, 220, GOLD, "1", "GROUNDED", "Answers begin with published Drupal content and retain source links."),
        (620, 378, SLATE, "2", "GOVERNED", "Access and future licensing rules remain under repository control."),
        (620, 536, GREEN, "3", "EXTENSIBLE", "New retrieval functions can be added through Drupal plugins."),
    ]
    for x, y, color, number, heading, body in principles:
        add_box(slide, x, y, 906, 132, WHITE, "#D1DADF")
        add_circle(slide, x + 28, y + 33, 66, color)
        add_text(slide, x + 28, y + 33, 66, 66, number, 17, INK if color == GOLD else WHITE, True, "center", "middle")
        add_text(slide, x + 122, y + 26, 250, 30, heading, 14, INK, True)
        add_text(slide, x + 122, y + 61, 735, 54, body, 13, SLATE)
    add_box(slide, 74, 744, 1452, 70, "#F1E7CA")
    add_text(slide, 100, 763, 170, 30, "DESIGN PRINCIPLE", 11, GOLD_DARK, True)
    add_text(slide, 300, 761, 1120, 34, "Drupal remains the source of truth; MCP is the access layer, not a replacement repository.", 14, INK, True)
    add_footer(slide, 2, "")

    # Slide 3 — what the prototype proves.
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(PAPER)
    add_title(slide, "02 • prototype", "Prototype proof", "What the current code demonstrates end to end.")
    stages = [
        (74, 220, 230, "AI CLIENT", "MCP request"),
        (348, 220, 250, "TRANSPORT", "HTTP or stdio"),
        (642, 220, 260, "DISPATCH", "JSON-RPC methods"),
        (946, 220, 260, "PLUGIN LAYER", "tools + resources"),
        (1250, 220, 276, "DRUPAL DATA", "nodes + PDFs"),
    ]
    for i, (x, y, width, heading, body) in enumerate(stages):
        fill = INK if i == 2 else WHITE
        add_box(slide, x, y, width, 122, fill, None if fill == INK else "#C9D2D8")
        add_text(slide, x + 10, y + 24, width - 20, 32, heading, 13, GOLD if fill == INK else INK, True, "center")
        add_text(slide, x + 10, y + 70, width - 20, 28, body, 11, WHITE if fill == INK else SLATE, False, "center")
        if i < len(stages) - 1:
            add_arrow(slide, x + width + 10, y + 49, stages[i + 1][0] - (x + width) - 20)
    capabilities = [
        (74, GOLD, "CONNECT", "2 transports", "Browser/server clients and local command-line clients."),
        (448, SLATE, "DISCOVER", "5 tools • 3 resources", "Clients can list capabilities before calling them."),
        (822, GREEN, "RETRIEVE", "nodes • fields • PDF page", "Published content is returned as structured JSON text."),
        (1196, INK_2, "EXTEND", "annotation plugins", "New functions are discovered through Drupal’s plugin system."),
    ]
    for x, color, heading, metric, body in capabilities:
        add_box(slide, x, 400, 330, 304, WHITE, "#D1DADF")
        add_box(slide, x, 400, 330, 12, color, rounded_box=False)
        add_text(slide, x + 28, 438, 250, 28, heading, 12, INK, True)
        add_text(slide, x + 28, 482, 280, 42, metric, 18, color, True)
        add_text(slide, x + 28, 548, 270, 88, body, 12, SLATE)
    add_box(slide, 74, 744, 1452, 70, INK)
    add_text(slide, 100, 763, 140, 30, "PROOF POINT", 11, GOLD, True)
    add_text(slide, 262, 761, 1160, 36, "A standard MCP client can discover and invoke Drupal functions without a repository-specific client.", 14, WHITE, True)
    add_footer(slide, 3, "")

    # Slide 4 — code-informed readiness gaps.
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(PAPER)
    add_title(slide, "03 • assessment", "Readiness gaps", "The prototype is extensible; access, operations and retrieval depth are not yet production-ready.")
    readiness_columns = [
        (74, RED, "ACCESS & EXPOSURE", [
            "Direct node/PDF reads do not enforce entity view access.",
            "One route permission covers every client and content type.",
            "Wildcard origins and responses expose email/file-path details.",
        ], "Risk: published but restricted content may be overexposed."),
        (574, GOLD, "RELIABILITY & OPERATIONS", [
            "HTTP and stdio duplicate the protocol dispatcher.",
            "A faulty plugin can disrupt tool/resource discovery.",
            "Logging/templates are partial; no behavioural test suite.",
        ], "Risk: maintenance drift and fragile client connections."),
        (1074, SLATE, "RETRIEVAL QUALITY", [
            "Search matches titles only; there is no ranking or semantic index.",
            "PDF extraction returns one text-based page at a time.",
            "No OCR, cached extraction or attachment-level search.",
        ], "Risk: the demo works, but evidence coverage is narrow."),
    ]
    for x, color, heading, items, implication in readiness_columns:
        add_box(slide, x, 220, 452, 494, WHITE, "#D1DADF")
        add_box(slide, x, 220, 452, 14, color, rounded_box=False)
        add_text(slide, x + 28, 258, 385, 32, heading, 13, INK, True)
        for i, item in enumerate(items):
            add_bullet(slide, x + 30, 316 + i * 94, item, color, width=350, size=12)
        add_box(slide, x + 24, 610, 404, 76, "#F4F6F7")
        add_text(slide, x + 42, 624, 365, 48, implication, 11, INK, True)
    add_box(slide, 74, 744, 1452, 70, INK)
    add_text(slide, 100, 763, 120, 30, "READINESS", 11, GOLD, True)
    add_text(slide, 244, 761, 1120, 36, "Appropriate for controlled internal testing only; access hardening is the release gate.", 14, WHITE, True)
    add_footer(slide, 4, "")

    # Slide 5 — pilot blueprint.
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(PAPER)
    add_title(slide, "04 • delivery", "Pilot blueprint", "Sequence the work around secure access, retrieval quality and operational assurance.")
    workstreams = [
        (74, GOLD, "SPRINT 4", "HARDEN ACCESS", [
            "Share one dispatcher and response model across both transports.",
            "Apply entity-level permissions and minimise response data.",
            "Add origin controls, usage limits and audit logging.",
        ], "Unauthorised content is never returned."),
        (574, SLATE, "SPRINT 5", "IMPROVE RETRIEVAL", [
            "Create a reusable extraction service with OCR and caching.",
            "Establish a search baseline, then prototype RAG.",
            "Return citation metadata and stable source links.",
        ], "Representative nodes and PDFs are retrievable with provenance."),
        (1074, GREEN, "SPRINT 6", "PROVE OPERATION", [
            "Test the same behaviours over HTTP and command line.",
            "Isolate plugin failures and make errors observable.",
            "Deploy to the test site and complete security review.",
        ], "The demonstration is repeatable and failures are diagnosable."),
    ]
    for x, color, sprint, heading, items, success in workstreams:
        add_box(slide, x, 220, 452, 494, WHITE, "#D1DADF")
        add_box(slide, x, 220, 452, 14, color, rounded_box=False)
        add_box(slide, x + 28, 260, 114, 36, "#EEF1F3")
        add_text(slide, x + 28, 260, 114, 36, sprint, 10, INK, True, "center", "middle")
        add_text(slide, x + 28, 322, 365, 32, heading, 14, INK, True)
        for i, item in enumerate(items):
            add_bullet(slide, x + 30, 372 + i * 78, item, color, width=350, size=11)
        add_box(slide, x + 24, 612, 404, 74, "#F4F6F7")
        add_text(slide, x + 42, 624, 80, 20, "SUCCESS", 9, color, True)
        add_text(slide, x + 42, 646, 360, 34, success, 10, INK, True)
    add_box(slide, 74, 744, 1452, 70, INK)
    add_text(slide, 100, 763, 150, 30, "FIRST DECISION", 11, GOLD, True)
    add_text(slide, 284, 761, 1120, 36, "Agree the pilot users and representative content set before expanding search features.", 14, WHITE, True)
    add_footer(slide, 5, "")

    prs.core_properties.title = "Policy Evidence Interface — Project Update"
    prs.core_properties.subject = "Project vision, Drupal MCP prototype proof, readiness assessment and pilot blueprint"
    prs.core_properties.author = "ANU TechLauncher Policy Evidence Interface Team"
    prs.core_properties.keywords = "Drupal, MCP, policy evidence, AI, TechLauncher"
    prs.save(OUTPUT)


def main():
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    slides = [simple_slide_1(), simple_slide_reason(), simple_slide_2(), simple_slide_3(), simple_slide_4()]
    paths: list[Path] = []
    for i, slide in enumerate(slides, 1):
        path = ASSET_DIR / f"slide_{i}.png"
        slide.save(path, quality=95)
        paths.append(path)
    make_contact_sheet(paths)
    build_pptx(paths)
    print(OUTPUT)


if __name__ == "__main__":
    main()
